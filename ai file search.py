import os
import re
import tkinter as tk
from tkinter import filedialog, messagebox, scrolledtext
from datetime import datetime

# FAISS and embedding libraries
import faiss
import numpy as np
from sentence_transformers import SentenceTransformer

# Document extraction libraries
import docx
import openpyxl
import pptx

# Ollama for LLM interaction
from ollama import chat

# Model configuration
MODEL = "gemma3:1b"
EMBEDDING_MODEL_NAME = "all-MiniLM-L6-v2"  # Change if needed

# Initialize the SentenceTransformer model for embeddings
embedder = SentenceTransformer(EMBEDDING_MODEL_NAME)

# ---------------------------
# LLM & Extraction Functions
# ---------------------------
def call_ollama(prompt: str) -> str:
    """Calls the Ollama model using the streaming chat interface and returns the response."""
    messages = [{'role': 'user', 'content': prompt}]
    try:
        stream = chat(model=MODEL, messages=messages, stream=True)
    except Exception as e:
        log(f"Error calling Ollama: {e}")
        return ""
    response = ""
    for chunk in stream:
        response += chunk['message']['content']
    return response.strip()

def extract_text_from_txt(filepath: str) -> str:
    try:
        with open(filepath, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        log(f"Error reading text file {filepath}: {e}")
        return ""

def extract_text_from_docx(filepath: str) -> str:
    try:
        doc = docx.Document(filepath)
        full_text = [para.text for para in doc.paragraphs]
        return "\n".join(full_text)
    except Exception as e:
        log(f"Error reading DOCX file {filepath}: {e}")
        return ""

def extract_text_from_xlsx(filepath: str) -> str:
    try:
        wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
        text_parts = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join([str(cell) for cell in row if cell is not None])
                text_parts.append(row_text)
        return "\n".join(text_parts)
    except Exception as e:
        log(f"Error reading XLSX file {filepath}: {e}")
        return ""

def extract_text_from_pptx(filepath: str) -> str:
    try:
        presentation = pptx.Presentation(filepath)
        text_runs = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    except Exception as e:
        log(f"Error reading PPTX file {filepath}: {e}")
        return ""

def extract_text(filepath: str) -> str:
    """Determines file type by extension and extracts text accordingly."""
    ext = os.path.splitext(filepath)[1].lower()
    if ext == ".txt":
        return extract_text_from_txt(filepath)
    elif ext == ".docx":
        return extract_text_from_docx(filepath)
    elif ext == ".xlsx":
        return extract_text_from_xlsx(filepath)
    elif ext == ".pptx":
        return extract_text_from_pptx(filepath)
    else:
        log(f"Unsupported file type: {filepath}")
        return ""

def summarize_document(filepath: str) -> str:
    """Extracts text from the document and generates a summary using the LLM."""
    content = extract_text(filepath)
    if not content.strip():
        log(f"No content extracted from {filepath}")
        return ""
    prompt = (f"Summarize the following document into key essential information that helps in search:\n\n"
              f"{content}\n\nSummary:")
    summary = call_ollama(prompt)
    return summary

def get_relevance_score(summary: str, query: str) -> float:
    """
    Uses the LLM to rate the relevance of the document summary with respect to the query.
    Returns a float score.
    """
    prompt = (f"On a scale from 1 to 10, where 10 means extremely relevant and 1 means not relevant, "
              f"rate the relevance of the following document summary to this query.\n\n"
              f"Query: {query}\n\nDocument Summary:\n{summary}\n\nRelevance rating (just the number):")
    score_str = call_ollama(prompt)
    try:
        if "/" in score_str:
            numerator = score_str.split("/")[0].strip()
            score = float(numerator)
        else:
            score = float(score_str.split()[0])
    except Exception as e:
        log(f"Error parsing relevance score from '{score_str}': {e}")
        score = 0.0
    return score

# ---------------------------
# Indexing & FAISS Functions
# ---------------------------
def index_documents(directory: str, extensions: list = None):
    """
    Walk through the directory, generate/update summaries for supported document types,
    and save each summary as a sidecar file (<filename>.summary.txt).
    """
    if extensions is None:
        extensions = [".txt", ".docx", ".xlsx", ".pptx"]
    count = 0
    for root, _, files in os.walk(directory):
        for file in files:
            if any(file.lower().endswith(ext) for ext in extensions):
                full_path = os.path.join(root, file)
                summary_path = full_path + ".summary.txt"
                if os.path.exists(summary_path):
                    doc_mtime = os.path.getmtime(full_path)
                    summary_mtime = os.path.getmtime(summary_path)
                    if summary_mtime >= doc_mtime:
                        log(f"Skipping {full_path}, summary is up-to-date.")
                        continue
                log(f"Indexing {full_path} ...")
                summary = summarize_document(full_path)
                if summary:
                    try:
                        with open(summary_path, 'w', encoding='utf-8') as sf:
                            sf.write(summary)
                        log(f"Summary saved to {summary_path}")
                        count += 1
                    except Exception as e:
                        log(f"Error writing summary for {full_path}: {e}")
    log(f"Indexing complete. {count} documents processed.")

def build_vector_index(directory: str, extensions: list = None):
    """
    Walk through the directory and load all summaries.
    Compute embeddings for each summary and build a FAISS index.
    Returns the FAISS index and a list of metadata dictionaries for each document.
    """
    if extensions is None:
        extensions = [".txt", ".docx", ".xlsx", ".pptx"]
    summaries = []
    meta_data = []  # Each dict contains 'doc_path', 'title', 'modified', and 'summary'
    for root, _, files in os.walk(directory):
        for file in files:
            for ext in extensions:
                if file.lower().endswith(ext + ".summary.txt"):
                    summary_path = os.path.join(root, file)
                    doc_path = summary_path[:-len(".summary.txt")]
                    if not os.path.exists(doc_path):
                        continue
                    try:
                        with open(summary_path, 'r', encoding='utf-8') as sf:
                            summary = sf.read()
                    except Exception as e:
                        log(f"Error reading summary {summary_path}: {e}")
                        continue
                    summaries.append(summary)
                    meta_data.append({
                        "doc_path": doc_path,
                        "title": os.path.basename(doc_path),
                        "modified": os.path.getmtime(doc_path),
                        "summary": summary
                    })
                    break
    if not summaries:
        return None, []
    embeddings = embedder.encode(summaries, convert_to_tensor=False)
    dimension = len(embeddings[0])
    index = faiss.IndexFlatL2(dimension)
    index.add(np.array(embeddings, dtype='float32'))
    return index, meta_data

# ---------------------------
# RAG-based Re-ranking Functions
# ---------------------------
def rerank_candidates_with_llm(candidate_docs: list, query: str) -> list:
    """
    Given a list of 5 candidate documents with their details, build a single prompt that includes:
      - The details of each candidate in the format:
        "Rank X: [Title] [date modified] [summary]"
      - The user query.
    The prompt instructs the LLM to return only the updated ranks as a comma-separated list (e.g., 3,1,4,2,5)
    with no extra details.
    If exactly 5 distinct numbers are returned, reorder the candidates accordingly.
    Otherwise, fall back to the original order.
    """
    prompt_lines = []
    prompt_lines.append("Below are 5 candidates shortlisted for top relevance. You should read the content below and rank accordingly to the user requirement:")
    for i, candidate in enumerate(candidate_docs, start=1):
        mod_time_str = datetime.fromtimestamp(candidate["modified"]).strftime('%Y-%m-%d %H:%M:%S')
        prompt_lines.append(f"Rank {i}: {candidate['title']} {mod_time_str} {candidate['summary']}")
    prompt_lines.append(f"\nUser Query: '{query}'")
    prompt_lines.append("Based on the user query, please return just the updated ranks in a comma-separated list (e.g., 3,1,4,2,5) with no extra details, where the first rank is the most relevant.")
    prompt = "\n".join(prompt_lines)
    llm_response = call_ollama(prompt)
    log(f"LLM re-ranking response: {llm_response}")
    numbers = re.findall(r'\d+', llm_response)
    unique_numbers = []
    for num_str in numbers:
        try:
            num = int(num_str)
            if 1 <= num <= len(candidate_docs) and num not in unique_numbers:
                unique_numbers.append(num)
        except Exception as e:
            log(f"Error parsing number from re-ranking response: {e}")
    if len(unique_numbers) == 5:
        reordered = [candidate_docs[num - 1] for num in unique_numbers]
        return reordered
    else:
        log("LLM did not return exactly 5 distinct numbers. Falling back to original order.")
        return candidate_docs

def search_documents_rag(directory: str, query: str, top_k: int = 10, final_k: int = 5, extensions: list = None) -> (list, list):
    """
    Uses a RAG setup to search documents:
      1. Build a FAISS vector index from summaries.
      2. Retrieve the top_k candidate documents.
      3. From these, select the top final_k candidates.
      4. Return both the initial candidate list (in retrieval order) and the re-ranked candidate list.
    """
    index, meta_data = build_vector_index(directory, extensions)
    if index is None:
        log("No summaries found in the directory.")
        return [], []
    query_emb = embedder.encode([query], convert_to_tensor=False)
    query_emb = np.array(query_emb, dtype='float32')
    distances, indices = index.search(query_emb, top_k)
    candidate_docs = []
    for idx in indices[0]:
        if idx < len(meta_data):
            candidate_docs.append(meta_data[idx])
    candidate_docs = candidate_docs[:final_k]
    initial_candidates = candidate_docs.copy()
    if not candidate_docs or len(candidate_docs) < final_k:
        log("Not enough candidates for re-ranking; using initial order.")
        return initial_candidates, initial_candidates
    reranked_candidates = rerank_candidates_with_llm(candidate_docs, query)
    return initial_candidates, reranked_candidates

# ---------------------------
# GUI Section
# ---------------------------
def log(message: str):
    """Append a message to the log text widget."""
    log_text.configure(state='normal')
    log_text.insert(tk.END, message + "\n")
    log_text.configure(state='disabled')
    log_text.see(tk.END)

def select_directory():
    dirname = filedialog.askdirectory()
    if dirname:
        dir_entry.delete(0, tk.END)
        dir_entry.insert(0, dirname)

def do_index():
    directory = dir_entry.get().strip()
    if not directory or not os.path.isdir(directory):
        messagebox.showerror("Error", "Please select a valid directory!")
        return
    log("Starting indexing...")
    index_documents(directory)
    log("Indexing finished.")

def do_search():
    directory = dir_entry.get().strip()
    query = query_entry.get().strip()
    if not directory or not os.path.isdir(directory):
        messagebox.showerror("Error", "Please select a valid directory!")
        return
    if not query:
        messagebox.showerror("Error", "Please enter a search query!")
        return
    log("Starting RAG-based search...")
    initial_candidates, reranked_candidates = search_documents_rag(directory, query)
    result_listbox.delete(0, tk.END)
    if not initial_candidates:
        result_listbox.insert(tk.END, "No relevant documents found.")
    else:
        result_listbox.insert(tk.END, "Initial Ranking:")
        for i, doc in enumerate(initial_candidates, start=1):
            mod_time_str = datetime.fromtimestamp(doc["modified"]).strftime('%Y-%m-%d %H:%M:%S')
            result_listbox.insert(tk.END, f"Rank {i}: {doc['title']} | Modified: {mod_time_str}")
        result_listbox.insert(tk.END, "---------------------")
        result_listbox.insert(tk.END, "Re-ranked Results:")
        if not reranked_candidates:
            result_listbox.insert(tk.END, "No re-ranked results found.")
        else:
            for i, doc in enumerate(reranked_candidates, start=1):
                mod_time_str = datetime.fromtimestamp(doc["modified"]).strftime('%Y-%m-%d %H:%M:%S')
                result_listbox.insert(tk.END, f"Rank {i}: {doc['title']} | Modified: {mod_time_str}")
    log("Search finished.")

# Create the main application window
root = tk.Tk()
root.title("RAG Document Index and Search")

# Frame for directory selection and indexing
dir_frame = tk.Frame(root)
dir_frame.pack(fill=tk.X, padx=10, pady=5)
tk.Label(dir_frame, text="Directory:").pack(side=tk.LEFT)
dir_entry = tk.Entry(dir_frame, width=50)
dir_entry.pack(side=tk.LEFT, padx=5)
tk.Button(dir_frame, text="Browse", command=select_directory).pack(side=tk.LEFT)
tk.Button(dir_frame, text="Index Documents", command=do_index).pack(side=tk.LEFT, padx=5)

# Frame for search query
query_frame = tk.Frame(root)
query_frame.pack(fill=tk.X, padx=10, pady=5)
tk.Label(query_frame, text="Search Query:").pack(side=tk.LEFT)
query_entry = tk.Entry(query_frame, width=50)
query_entry.pack(side=tk.LEFT, padx=5)
tk.Button(query_frame, text="Search Documents", command=do_search).pack(side=tk.LEFT, padx=5)

# Frame for displaying search results
result_frame = tk.Frame(root)
result_frame.pack(fill=tk.BOTH, expand=True, padx=10, pady=5)
tk.Label(result_frame, text="Search Results:").pack(anchor=tk.W)
result_listbox = tk.Listbox(result_frame, height=15)
result_listbox.pack(fill=tk.BOTH, expand=True)

# Log window (for status and debugging)
log_frame = tk.Frame(root)
log_frame.pack(fill=tk.BOTH, expand=True, padx=10, pady=5)
tk.Label(log_frame, text="Log:").pack(anchor=tk.W)
log_text = scrolledtext.ScrolledText(log_frame, height=10, state='disabled')
log_text.pack(fill=tk.BOTH, expand=True)

root.mainloop()
